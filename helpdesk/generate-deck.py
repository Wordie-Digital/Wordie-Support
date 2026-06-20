"""
Wordie Support Operations — deck generator
Produces helpdesk/Wordie-Support-Workflows.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand tokens ──────────────────────────────────────────────────────────────
DARK_TEAL   = RGBColor(0x0A, 0x35, 0x42)
ACCENT_TEAL = RGBColor(0x11, 0x6E, 0x6E)
LIGHT_GREEN = RGBColor(0xB9, 0xDD, 0xDD)
CORAL       = RGBColor(0xF5, 0x63, 0x4D)
NEAR_BLACK  = RGBColor(0x06, 0x20, 0x28)
OFF_WHITE   = RGBColor(0xF6, 0xF9, 0xF9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEAD = "Montserrat"
FONT_BODY = "Source Sans Pro"

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def label(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE,
          font=FONT_BODY, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def divider(slide, y, color=ACCENT_TEAL, thickness=0.03):
    shape = slide.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(12.13), Inches(thickness))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def pill(slide, text, x, y, w=1.6, h=0.38, bg=CORAL, text_color=WHITE, size=13):
    r = rect(slide, x, y, w, h, bg)
    r.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name  = FONT_HEAD
    run.font.size  = Pt(size)
    run.font.bold  = True
    run.font.color.rgb = text_color


def section_header(slide, number, title):
    """Dark teal band across top with section number + title."""
    rect(slide, 0, 0, 13.33, 1.4, DARK_TEAL)
    label(slide, number, 0.55, 0.28, 1.0, 0.9,
          size=36, bold=True, color=CORAL, font=FONT_HEAD)
    label(slide, title,  1.45, 0.35, 10.0, 0.75,
          size=26, bold=True, color=WHITE, font=FONT_HEAD)


def footer(slide, text="Wordie Digital  ·  support@wordie.com.au  ·  wordie.com.au"):
    rect(slide, 0, 7.18, 13.33, 0.32, NEAR_BLACK)
    label(slide, text, 0.55, 7.2, 12.0, 0.3,
          size=9, color=LIGHT_GREEN, font=FONT_BODY)


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, DARK_TEAL)
    # diagonal accent ribbon
    from pptx.util import Pt as _Pt
    accent = s.shapes.add_shape(1, Inches(8.5), Inches(-0.5), Inches(6), Inches(9))
    accent.rotation = -18
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_TEAL
    accent.line.fill.background()
    accent2 = s.shapes.add_shape(1, Inches(10.2), Inches(-0.5), Inches(3), Inches(9))
    accent2.rotation = -18
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = RGBColor(0x0D, 0x55, 0x55)
    accent2.line.fill.background()

    label(s, "Support Operations", 0.7, 1.8, 8.0, 0.7,
          size=16, color=LIGHT_GREEN, font=FONT_BODY, italic=True)
    label(s, "How our AI-assisted\nsupport works", 0.7, 2.3, 9.0, 2.4,
          size=44, bold=True, color=WHITE, font=FONT_HEAD)
    label(s, "Wordie Digital  ·  wordie.com.au", 0.7, 6.8, 6.0, 0.4,
          size=11, color=LIGHT_GREEN, font=FONT_BODY)


def slide_overview(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
    section_header(s, "01", "Three systems, working together")
    footer(s)

    items = [
        (DARK_TEAL,   "HubSpot Auto-Triage",        "Every hour, Claude classifies new tickets,\nadds structured notes, and sets priority\nand category fields automatically."),
        (ACCENT_TEAL, "Helpdesk Knowledge Base",    "Canonical playbooks for recurring issues —\ndiagnosis steps, fix instructions, and\ncopy-paste agent scripts."),
        (CORAL,       "GitHub Deploy Workflow",      "All code changes go through GitHub Actions\nbefore reaching WP Engine — never a direct push."),
    ]

    for i, (color, heading, body) in enumerate(items):
        x = 0.55 + i * 4.2
        rect(s, x, 1.7, 3.8, 4.6, color)
        label(s, heading, x + 0.25, 1.95, 3.3, 0.7,
              size=17, bold=True, color=WHITE, font=FONT_HEAD)
        label(s, body, x + 0.25, 2.75, 3.3, 2.8,
              size=13, color=WHITE, font=FONT_BODY)

    label(s, "Each system feeds the next. Triage notes reference playbooks. Playbooks are version-controlled in GitHub.",
          0.55, 6.55, 12.2, 0.45, size=11, color=ACCENT_TEAL, font=FONT_BODY, italic=True)


def slide_triage_how(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
    section_header(s, "02", "Auto-triage — how it runs")
    footer(s)

    steps = [
        ("1", DARK_TEAL,   "Every hour\nat :00",        "Claude Code wakes up and reads\nthe task definition file."),
        ("2", ACCENT_TEAL, "Fetches\nnew tickets",      "Searches HubSpot for open tickets\ncreated in the last 2 hours."),
        ("3", ACCENT_TEAL, "Skips\nalready triaged",    "Any ticket with an AUTO-TRIAGE\nnote is ignored — no duplicates."),
        ("4", DARK_TEAL,   "Classifies\nP1 – P4",       "Applies priority logic and selects\na routing destination."),
        ("5", CORAL,       "Writes note\nto HubSpot",   "Posts a structured triage note\ndirectly onto the ticket."),
        ("6", CORAL,       "Updates\nfields",           "Sets Priority + Category fields\nif currently blank."),
    ]

    arrow_color = LIGHT_GREEN
    for i, (num, color, heading, body) in enumerate(steps):
        x = 0.4 + i * 2.1
        rect(s, x, 1.65, 1.85, 4.8, color)
        label(s, num,     x+0.15, 1.82, 0.5, 0.5, size=22, bold=True, color=WHITE, font=FONT_HEAD)
        label(s, heading, x+0.12, 2.4,  1.6, 0.9, size=13, bold=True, color=WHITE, font=FONT_HEAD)
        label(s, body,    x+0.12, 3.45, 1.6, 2.0, size=11, color=WHITE, font=FONT_BODY)
        if i < len(steps) - 1:
            ax = x + 1.85
            label(s, "→", ax, 3.55, 0.28, 0.4, size=18, bold=True, color=ACCENT_TEAL, font=FONT_HEAD)

    label(s, "Skips: M&S retainer tickets · Pure project delivery · Inbound call records · Government notifications",
          0.55, 6.55, 12.2, 0.45, size=11, color=ACCENT_TEAL, font=FONT_BODY, italic=True)


def slide_triage_note(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
    section_header(s, "02", "Auto-triage — what the note looks like")
    footer(s)

    # Note preview panel
    rect(s, 0.55, 1.65, 6.2, 5.1, WHITE)
    # top bar
    rect(s, 0.55, 1.65, 6.2, 0.45, DARK_TEAL)
    label(s, "HubSpot ticket note — added automatically", 0.75, 1.7, 5.8, 0.38,
          size=11, color=LIGHT_GREEN, font=FONT_BODY, italic=True)

    note_lines = [
        ("🔍 AUTO-TRIAGE  |  2026-05-19", 13, True,  DARK_TEAL),
        ("", 6, False, NEAR_BLACK),
        ("INTENT:  Client unable to upload product images in WordPress", 11, False, NEAR_BLACK),
        ("PRIORITY:  P3  |  URGENCY:  normal", 11, False, NEAR_BLACK),
        ("ROUTED TO:  Developer  |  ESCALATION:  No", 11, False, NEAR_BLACK),
        ("", 6, False, NEAR_BLACK),
        ("SUMMARY:  The uploads directory has lost write permissions,", 11, False, NEAR_BLACK),
        ("preventing any media uploads. No revenue impact but blocks", 11, False, NEAR_BLACK),
        ("content publishing workflow.", 11, False, NEAR_BLACK),
        ("", 6, False, NEAR_BLACK),
        ("RECOMMENDED ACTION:  SSH/SFTP to correct permissions on", 11, False, NEAR_BLACK),
        ("wp-content/uploads. See playbook: media-upload-error.md", 11, False, NEAR_BLACK),
        ("", 6, False, NEAR_BLACK),
        ("REUSABILITY:  8/10  |  KNOWN PATTERN ✓", 11, True, ACCENT_TEAL),
        ("TAGS:  wordpress, uploads, permissions, p3", 11, False, NEAR_BLACK),
    ]

    y = 2.22
    for line, size, bold, color in note_lines:
        label(s, line, 0.72, y, 5.8, 0.28, size=size, bold=bold, color=color, font=FONT_BODY)
        y += size * 0.018 + 0.01

    # Callouts on the right
    callouts = [
        (CORAL,       "INTENT",             "Plain-English description of what the client needs."),
        (DARK_TEAL,   "PRIORITY + ROUTING", "P1-P4 severity and which team handles it."),
        (ACCENT_TEAL, "RECOMMENDED ACTION", "The exact next step — no guesswork for your team."),
        (DARK_TEAL,   "KNOWN PATTERN ✓",    "Matched to a helpdesk playbook. Follow that doc for the fix."),
    ]

    y = 1.7
    for color, heading, body in callouts:
        rect(s, 7.1, y, 0.08, 0.78, color)
        label(s, heading, 7.3, y,       5.6, 0.32, size=12, bold=True, color=color, font=FONT_HEAD)
        label(s, body,    7.3, y+0.3,   5.6, 0.55, size=11, color=NEAR_BLACK, font=FONT_BODY)
        y += 1.1


def slide_priority(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
    section_header(s, "03", "Priority framework")
    footer(s)

    rows = [
        (CORAL,       "P1  URGENT",  "< 1 hour",   "Site down · Checkout broken · Security breach · Emails completely stopped"),
        (DARK_TEAL,   "P2  HIGH",    "< 4 hours",  "Major feature broken · Login blocked · Analytics gone · ERP sync down · Backups failing"),
        (ACCENT_TEAL, "P3  MEDIUM",  "< 24 hours", "Partial issue with a workaround · Styling bugs · Upload permissions · 404 errors"),
        (LIGHT_GREEN, "P4  LOW",     "< 72 hours", "Cosmetic · Content update · Meta tags · Redirects · General question"),
    ]

    for i, (color, label_text, sla, description) in enumerate(rows):
        y = 1.65 + i * 1.3
        rect(s, 0.55, y, 2.5, 1.1, color)
        tc = WHITE if color != LIGHT_GREEN else DARK_TEAL
        label(s, label_text,   0.65, y+0.18, 2.3, 0.5, size=15, bold=True, color=tc, font=FONT_HEAD)
        label(s, "SLA  " + sla, 0.65, y+0.65, 2.3, 0.35, size=11, color=tc, font=FONT_BODY)
        label(s, description,  3.25, y+0.28, 9.5, 0.7, size=13, color=NEAR_BLACK, font=FONT_BODY)
        divider(s, y + 1.15, color=LIGHT_GREEN, thickness=0.02)


def slide_routing(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
    section_header(s, "04", "Routing guide")
    footer(s)

    routes = [
        (CORAL,       "Developer",        "Code bugs, JS errors, integrations, analytics setup,\nCSS requiring code changes, performance issues."),
        (DARK_TEAL,   "Support Agent",    "Content updates, IP whitelists, meta tags, redirects,\naccess requests, test accounts — you can close these."),
        (ACCENT_TEAL, "Escalation Team",  "Security issues, data loss, active breaches,\npayment system down. Flag to Stewart immediately."),
        (LIGHT_GREEN, "Chatbot",          "Spam, misdirected email, government notifications, FAQs.\nClose or send the self-serve template."),
    ]

    for i, (color, heading, body) in enumerate(routes):
        col = i % 2
        row = i // 2
        x = 0.55 + col * 6.4
        y = 1.65 + row * 2.7
        rect(s, x, y, 6.0, 2.45, color)
        tc = WHITE if color != LIGHT_GREEN else DARK_TEAL
        label(s, heading, x+0.25, y+0.25, 5.5, 0.55, size=18, bold=True, color=tc, font=FONT_HEAD)
        label(s, body,    x+0.25, y+0.9,  5.5, 1.3,  size=13, color=tc, font=FONT_BODY)


def slide_playbooks(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
    section_header(s, "05", "Helpdesk knowledge base")
    footer(s)

    label(s, "Each playbook covers: symptom · root cause · diagnosis · fix · prevention · agent script · chatbot response",
          0.55, 1.5, 12.2, 0.42, size=12, color=ACCENT_TEAL, font=FONT_BODY, italic=True)

    playbooks = [
        ("WooCommerce order flow failure",  "woocommerce/order-flow-failure.md",     "9/10",  "WP-Cron cleared → stuck orders / missing emails"),
        ("WordPress media upload error",    "wordpress/media-upload-error.md",       "8/10",  "Upload directory not writable — permissions fix"),
        ("Member login blocked",            "wordpress/member-login-blocked.md",     "8/10",  "Wordfence IP block — 3 client occurrences in 2026"),
        ("GTM container replaced",          "analytics/gtm-container-replaced.md",  "10/10", "GTM ID swapped on deploy — GTM-WL7TVVS is correct"),
        ("Debug code on production",        "wordpress/debug-code-on-production.md", "9/10",  "WP_DEBUG left on — treat as P1 immediately"),
    ]

    for i, (name, path, confidence, desc) in enumerate(playbooks):
        y = 2.1 + i * 0.98
        rect(s, 0.55, y, 0.08, 0.72, ACCENT_TEAL)
        label(s, name,   0.8, y+0.04, 4.5, 0.38, size=14, bold=True, color=DARK_TEAL, font=FONT_HEAD)
        label(s, desc,   0.8, y+0.42, 6.0, 0.35, size=11, color=NEAR_BLACK, font=FONT_BODY)
        # confidence pill
        conf_color = CORAL if confidence == "10/10" else ACCENT_TEAL
        rect(s, 7.1,  y+0.15, 1.0, 0.38, conf_color)
        label(s, confidence, 7.1, y+0.15, 1.0, 0.38, size=12, bold=True, color=WHITE, font=FONT_HEAD, align=PP_ALIGN.CENTER)
        label(s, path, 8.35, y+0.22, 4.8, 0.3, size=10, color=ACCENT_TEAL, font=FONT_BODY, italic=True)


def slide_scheduled_task(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, DARK_TEAL)
    footer(s)

    label(s, "06", 0.55, 0.3, 1.2, 0.9, size=36, bold=True, color=CORAL, font=FONT_HEAD)
    label(s, "What is a Claude Code scheduled task?", 1.45, 0.38, 10.5, 0.7,
          size=26, bold=True, color=WHITE, font=FONT_HEAD)
    divider(s, 1.25, color=CORAL)

    label(s,
          "Claude Code is Anthropic's AI coding assistant. Its Scheduled Tasks feature lets you define\n"
          "an AI agent in plain language and run it on a timer — like a cron job, but using natural\n"
          "language instructions instead of code.",
          0.55, 1.45, 12.2, 1.1, size=14, color=LIGHT_GREEN, font=FONT_BODY)

    # two columns
    col1 = [
        ("Task file",    "SKILL.md at\n~/.claude/scheduled-tasks/hubspot-auto-triage/"),
        ("Schedule",     "0 * * * *\n— every hour on the hour"),
        ("Connection",   "HubSpot MCP (secure API bridge)\n— no passwords stored in the script"),
    ]
    col2 = [
        ("Where to find it", "Claude Code desktop app\n→ clock icon in sidebar → Scheduled Tasks"),
        ("First run",        "Click Run now once to pre-approve\nHubSpot tool permissions"),
        ("If a run fails",   "Check run history in the app.\nUsually a timeout — just retry manually."),
    ]

    for i, (heading, body) in enumerate(col1):
        y = 2.75 + i * 1.38
        rect(s, 0.55, y, 0.06, 0.95, CORAL)
        label(s, heading, 0.75, y,      5.5, 0.38, size=13, bold=True, color=WHITE, font=FONT_HEAD)
        label(s, body,    0.75, y+0.42, 5.5, 0.65, size=11, color=LIGHT_GREEN, font=FONT_BODY)

    for i, (heading, body) in enumerate(col2):
        y = 2.75 + i * 1.38
        rect(s, 6.85, y, 0.06, 0.95, ACCENT_TEAL)
        label(s, heading, 7.05, y,      5.8, 0.38, size=13, bold=True, color=WHITE, font=FONT_HEAD)
        label(s, body,    7.05, y+0.42, 5.8, 0.65, size=11, color=LIGHT_GREEN, font=FONT_BODY)


def slide_deploy(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
    section_header(s, "07", "Deploy workflow — code to WP Engine")
    footer(s)

    steps = [
        (DARK_TEAL,   "Local\ndevelopment",    "Write and test code\non your machine."),
        (DARK_TEAL,   "git push\nto GitHub",   "Push your branch.\nNever push to WP Engine directly."),
        (ACCENT_TEAL, "Open Pull\nRequest",     "Request a review\nbefore anything goes live."),
        (ACCENT_TEAL, "Review\n+ merge",        "Approve and merge\ninto main."),
        (CORAL,       "GitHub Actions\nruns",   "Automated checks run\nbefore deploying."),
        (CORAL,       "WP Engine\ndeployment",  "Staging or production\nupdated automatically."),
    ]

    for i, (color, heading, body) in enumerate(steps):
        x = 0.45 + i * 2.1
        rect(s, x, 1.75, 1.85, 4.4, color)
        label(s, str(i+1), x+0.14, 1.92, 0.5, 0.5, size=20, bold=True, color=WHITE, font=FONT_HEAD)
        label(s, heading, x+0.12, 2.5,  1.62, 0.9, size=13, bold=True, color=WHITE, font=FONT_HEAD)
        label(s, body,    x+0.12, 3.55, 1.62, 1.8, size=11, color=WHITE, font=FONT_BODY)
        if i < len(steps) - 1:
            label(s, "→", x+1.85, 3.6, 0.3, 0.4, size=18, bold=True, color=ACCENT_TEAL, font=FONT_HEAD)

    label(s, "Rule: all changes go through GitHub → GitHub Actions. This prevents regressions like a GTM container ID being overwritten on a direct push.",
          0.55, 6.42, 12.2, 0.55, size=11, color=ACCENT_TEAL, font=FONT_BODY, italic=True)


def slide_quickref(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
    section_header(s, "08", "Quick reference for the team")
    footer(s)

    items = [
        (CORAL,       "Client can't log in",          "Check triage note. Usually Wordfence IP block → member-login-blocked.md"),
        (DARK_TEAL,   "Orders stuck / emails missing","P1 if checkout broken. WP-Cron was cleared → order-flow-failure.md"),
        (ACCENT_TEAL, "Can't upload images",          "Directory permissions fix — 15 min job → media-upload-error.md"),
        (DARK_TEAL,   "Analytics dropped to zero",    "Check GTM container ID in page source. Correct ID: GTM-WL7TVVS → gtm-container-replaced.md"),
        (CORAL,       "PHP errors on live site",      "P1 immediately. Debug mode left on → debug-code-on-production.md. Escalate to developer now."),
    ]

    for i, (color, trigger, action) in enumerate(items):
        y = 1.62 + i * 1.02
        rect(s, 0.55, y, 0.1, 0.82, color)
        label(s, trigger, 0.8, y+0.06,  4.8, 0.38, size=13, bold=True, color=DARK_TEAL, font=FONT_HEAD)
        label(s, action,  0.8, y+0.46,  11.9, 0.38, size=12, color=NEAR_BLACK, font=FONT_BODY)


def slide_closing(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, DARK_TEAL)
    # accent ribbons
    for ox, ow, alpha_color in [(8.2, 6, ACCENT_TEAL), (10.1, 3, RGBColor(0x0D, 0x55, 0x55))]:
        a = s.shapes.add_shape(1, Inches(ox), Inches(-0.5), Inches(ow), Inches(9))
        a.rotation = -18
        a.fill.solid()
        a.fill.fore_color.rgb = alpha_color
        a.line.fill.background()

    label(s, "Questions?", 0.7, 1.9, 8.0, 1.1,
          size=46, bold=True, color=WHITE, font=FONT_HEAD)
    label(s, "Ping Stewart on Slack, or reply to a HubSpot ticket thread.\nThe auto-triage system is maintained by Stewart — if a classification looks wrong, let him know.",
          0.7, 3.1, 8.5, 1.4, size=15, color=LIGHT_GREEN, font=FONT_BODY)
    label(s, "support@wordie.com.au", 0.7, 4.7, 6.0, 0.5,
          size=16, bold=True, color=CORAL, font=FONT_HEAD)
    label(s, "wordie.com.au", 0.7, 5.25, 4.0, 0.4,
          size=13, color=LIGHT_GREEN, font=FONT_BODY)


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_overview(prs)
    slide_triage_how(prs)
    slide_triage_note(prs)
    slide_priority(prs)
    slide_routing(prs)
    slide_playbooks(prs)
    slide_scheduled_task(prs)
    slide_deploy(prs)
    slide_quickref(prs)
    slide_closing(prs)

    out = "/Users/stewartlemalu/Documents/GitHub/.claude/worktrees/festive-saha-9dab7c/helpdesk/Wordie-Support-Workflows.pptx"
    prs.save(out)
    print(f"Saved → {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
